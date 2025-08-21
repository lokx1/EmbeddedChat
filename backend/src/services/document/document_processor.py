# Document processing service
import os
import base64
import asyncio
from typing import Dict, Any, Optional, List, Tuple
from pathlib import Path
import mimetypes
import logging

# PDF processing
try:
    import PyPDF2
    import fitz  # PyMuPDF for better PDF extraction
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Image processing
try:
    from PIL import Image
    import io
    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False

# Document processing
try:
    import docx
    import openpyxl
    import csv
    OFFICE_AVAILABLE = True
except ImportError:
    OFFICE_AVAILABLE = False

# PowerPoint processing
try:
    from pptx import Presentation
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False

# LaTeX processing
try:
    import pylatexenc.latex2text
    LATEX_AVAILABLE = True
except ImportError:
    LATEX_AVAILABLE = False

# Additional document formats - textract is optional due to system dependencies
ADVANCED_TEXT_AVAILABLE = False  # Disabled textract to avoid installation issues


from ..workflow.ai_providers import AIProviderFactory

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Service for processing and analyzing documents with AI"""
    
    def __init__(self):
        self.supported_image_types = {
            'image/jpeg', 'image/jpg', 'image/png', 'image/gif', 
            'image/webp', 'image/bmp', 'image/tiff'
        }
        self.supported_document_types = {
            'application/pdf', 'text/plain', 'text/markdown',
            'application/msword', 
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'text/csv', 'application/json',
            'application/vnd.ms-excel',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            # PowerPoint formats
            'application/vnd.ms-powerpoint',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            # LaTeX formats
            'application/x-latex', 'text/x-latex', 'application/x-tex', 'text/x-tex',
            # Rich text formats
            'application/rtf', 'text/rtf',
            # OpenDocument formats
            'application/vnd.oasis.opendocument.text',
            'application/vnd.oasis.opendocument.presentation',
            'application/vnd.oasis.opendocument.spreadsheet'
        }
        
        # Audio types supported by Gemini
        self.supported_audio_types = {
            'audio/wav', 'audio/mp3', 'audio/aiff', 'audio/aac', 'audio/ogg', 'audio/flac'
        }
        
        # Video types supported by Gemini
        self.supported_video_types = {
            'video/mp4', 'video/mpeg', 'video/mov', 'video/avi', 'video/x-flv', 
            'video/mpg', 'video/webm', 'video/wmv', 'video/3gpp'
        }
    
    def _clean_text_for_db(self, text: str) -> str:
        """Clean text to make it safe for PostgreSQL storage"""
        if not text:
            return ""
        
        # Remove null bytes and other problematic characters
        text = text.replace('\x00', '')  # Remove null bytes
        text = text.replace('\ufeff', '')  # Remove BOM
        
        # Remove other control characters that might cause issues
        import re
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    async def process_document(
        self, 
        file_path: str, 
        mime_type: str,
        ai_provider: str = "openai",
        api_key: str = None
    ) -> Dict[str, Any]:
        """
        Process a document and extract content/analysis
        
        Args:
            file_path: Path to the file
            mime_type: MIME type of the file
            ai_provider: AI provider to use for analysis
            api_key: API key for the AI provider
            
        Returns:
            Dict containing extracted content and AI analysis
        """
        try:
            result = {
                "success": False,
                "extracted_text": "",
                "analysis": "",
                "summary": "",
                "file_type": mime_type,
                "error": None
            }
            
            # Extract content based on file type
            if mime_type in self.supported_image_types:
                content_result = await self._process_image(file_path, ai_provider, api_key)
            elif mime_type in self.supported_document_types:
                content_result = await self._process_document_file(file_path, mime_type, ai_provider, api_key)
            elif mime_type in self.supported_video_types:
                content_result = await self._process_video(file_path, ai_provider, api_key)
            elif mime_type in self.supported_audio_types:
                content_result = await self._process_audio(file_path, ai_provider, api_key)
            else:
                result["error"] = f"Unsupported file type: {mime_type}"
                return result
            
            result.update(content_result)
            result["success"] = True
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            return {
                "success": False,
                "extracted_text": "",
                "analysis": "",
                "summary": "",
                "file_type": mime_type,
                "error": str(e)
            }
    
    async def _process_image(
        self, 
        file_path: str, 
        ai_provider: str, 
        api_key: str
    ) -> Dict[str, Any]:
        """Process image using AI vision capabilities"""
        try:
            # Read and encode image
            with open(file_path, 'rb') as f:
                image_data = f.read()
            
            # Convert to base64 for AI processing
            image_base64 = base64.b64encode(image_data).decode('utf-8')
            
            # Use AI provider for image analysis
            if ai_provider == "openai" and api_key:
                analysis = await self._analyze_image_with_openai(image_base64, api_key)
            elif ai_provider == "claude" and api_key:
                analysis = await self._analyze_image_with_claude(image_base64, api_key)
            elif ai_provider == "gemini" and api_key:
                analysis = await self._analyze_image_with_gemini(image_base64, api_key)
            else:
                # Basic image info without AI
                analysis = await self._get_basic_image_info(file_path)
            
            return {
                "extracted_text": self._clean_text_for_db(analysis.get("description", "")),
                "analysis": self._clean_text_for_db(analysis.get("analysis", "")),
                "summary": self._clean_text_for_db(analysis.get("summary", "")),
                "metadata": analysis.get("metadata", {})
            }
            
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {str(e)}")
            raise
    
    async def _process_document_file(
        self, 
        file_path: str, 
        mime_type: str,
        ai_provider: str,
        api_key: str
    ) -> Dict[str, Any]:
        """Process document files and extract text"""
        try:
            extracted_text = ""
            
            # Extract text based on file type
            if mime_type == "application/pdf":
                # For PDF, try Gemini direct processing first if AI provider is available
                if ai_provider == "gemini" and api_key:
                    try:
                        pdf_result = await self._process_pdf_with_gemini(file_path, ai_provider, api_key)
                        return pdf_result
                    except Exception as e:
                        logger.warning(f"Gemini PDF processing failed, falling back to text extraction: {str(e)}")
                        extracted_text = await self._extract_pdf_text(file_path)
                else:
                    extracted_text = await self._extract_pdf_text(file_path)
            elif mime_type == "text/plain" or mime_type == "text/markdown":
                extracted_text = await self._extract_text_file(file_path)
            elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                extracted_text = await self._extract_docx_text(file_path)
            elif mime_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
                extracted_text = await self._extract_ppt_text(file_path)
            elif mime_type in ["application/x-latex", "text/x-latex", "application/x-tex", "text/x-tex"]:
                extracted_text = await self._extract_latex_text(file_path)
            elif mime_type in ["application/rtf", "text/rtf"]:
                extracted_text = await self._extract_rtf_text(file_path)
            elif mime_type in ["application/vnd.oasis.opendocument.text", "application/vnd.oasis.opendocument.presentation"]:
                extracted_text = await self._extract_odt_text(file_path)
            elif mime_type == "text/csv":
                extracted_text = await self._extract_csv_text(file_path)
            elif mime_type == "application/json":
                extracted_text = await self._extract_json_text(file_path)
            else:
                extracted_text = "File content could not be extracted"
            
            # Generate AI analysis if provider is available
            analysis = ""
            summary = ""
            if extracted_text and ai_provider and api_key:
                ai_result = await self._analyze_text_with_ai(extracted_text, ai_provider, api_key)
                analysis = ai_result.get("analysis", "")
                summary = ai_result.get("summary", "")
            
            return {
                "extracted_text": self._clean_text_for_db(extracted_text),
                "analysis": self._clean_text_for_db(analysis),
                "summary": self._clean_text_for_db(summary),
                "metadata": {"text_length": len(extracted_text)}
            }
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            raise
    
    async def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            return "PDF processing libraries not available"
        
        try:
            text = ""
            
            # Try PyMuPDF first (better extraction)
            try:
                doc = fitz.open(file_path)
                for page in doc:
                    text += page.get_text()
                doc.close()
                return text
            except:
                pass
            
            # Fallback to PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text()
            
            return text
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            return "Error extracting PDF content"
    
    async def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            # Try to detect and handle BOM
            with open(file_path, 'rb') as file:
                raw_data = file.read()
            
            # Remove BOM if present
            if raw_data.startswith(b'\xff\xfe'):  # UTF-16 LE BOM
                text = raw_data[2:].decode('utf-16-le', errors='ignore')
            elif raw_data.startswith(b'\xfe\xff'):  # UTF-16 BE BOM  
                text = raw_data[2:].decode('utf-16-be', errors='ignore')
            elif raw_data.startswith(b'\xef\xbb\xbf'):  # UTF-8 BOM
                text = raw_data[3:].decode('utf-8', errors='ignore')
            else:
                # Try UTF-8 first
                try:
                    text = raw_data.decode('utf-8')
                except UnicodeDecodeError:
                    # Try different encodings
                    for encoding in ['latin-1', 'cp1252', 'iso-8859-1', 'utf-16']:
                        try:
                            text = raw_data.decode(encoding, errors='ignore')
                            break
                        except:
                            continue
                    else:
                        return "Error: Could not decode file with available encodings"
            
            # Clean up text for database storage
            return self._clean_text_for_db(text)
            
        except Exception as e:
            logger.error(f"Error extracting text file: {str(e)}")
            return "Error extracting text content"
    
    async def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX files"""
        if not OFFICE_AVAILABLE:
            return "Office document processing libraries not available"
        
        try:
            doc = docx.Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {str(e)}")
            return "Error extracting DOCX content"
    
    async def _extract_csv_text(self, file_path: str) -> str:
        """Extract text from CSV files"""
        try:
            # Read raw data first to handle encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
            
            # Decode with encoding detection
            try:
                # Remove BOM and decode
                if raw_data.startswith(b'\xef\xbb\xbf'):  # UTF-8 BOM
                    text_data = raw_data[3:].decode('utf-8', errors='ignore')
                elif raw_data.startswith(b'\xff\xfe'):  # UTF-16 LE BOM
                    text_data = raw_data[2:].decode('utf-16-le', errors='ignore')
                else:
                    text_data = raw_data.decode('utf-8', errors='ignore')
                    
                # Clean null bytes
                text_data = text_data.replace('\x00', '')
                
                # Parse CSV from string
                import io
                csv_reader = csv.reader(io.StringIO(text_data))
                rows = []
                for row in csv_reader:
                    rows.append(', '.join(row))
                return '\n'.join(rows)
                
            except Exception:
                # Fallback: try latin-1
                text_data = raw_data.decode('latin-1', errors='ignore')
                text_data = text_data.replace('\x00', '')
                csv_reader = csv.reader(io.StringIO(text_data))
                rows = []
                for row in csv_reader:
                    rows.append(', '.join(row))
                return '\n'.join(rows)
                
        except Exception as e:
            logger.error(f"Error extracting CSV text: {str(e)}")
            return "Error extracting CSV content"
    
    async def _extract_json_text(self, file_path: str) -> str:
        """Extract text from JSON files"""
        try:
            import json
            
            # Read raw data first to handle encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
            
            # Remove BOM and decode
            if raw_data.startswith(b'\xef\xbb\xbf'):  # UTF-8 BOM
                text_data = raw_data[3:].decode('utf-8', errors='ignore')
            elif raw_data.startswith(b'\xff\xfe'):  # UTF-16 LE BOM
                text_data = raw_data[2:].decode('utf-16-le', errors='ignore')
            else:
                text_data = raw_data.decode('utf-8', errors='ignore')
            
            # Clean null bytes
            text_data = text_data.replace('\x00', '')
            
            # Parse JSON
            data = json.loads(text_data)
            return json.dumps(data, indent=2, ensure_ascii=False)
            
        except Exception as e:
            logger.error(f"Error extracting JSON text: {str(e)}")
            return "Error extracting JSON content"
    
    async def _extract_ppt_text(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        try:
            if not PPT_AVAILABLE:
                return "PowerPoint processing library not available"
            
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_runs = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"=== Slide {slide_num} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.shape_type == 19:  # Table
                        try:
                            table = shape.table
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    row_text.append(cell.text.strip())
                                slide_text += " | ".join(row_text) + "\n"
                        except:
                            pass
                
                text_runs.append(slide_text)
            
            extracted_text = "\n".join(text_runs)
            return self._clean_text_for_db(extracted_text)
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint text: {str(e)}")
            return "Error extracting PowerPoint content"
    
    async def _extract_latex_text(self, file_path: str) -> str:
        """Extract text from LaTeX files"""
        try:
            if not LATEX_AVAILABLE:
                # Fallback to plain text extraction
                return await self._extract_text_file(file_path)
            
            from pylatexenc.latex2text import LatexNodes2Text
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                latex_content = file.read()
            
            # Convert LaTeX to plain text
            converter = LatexNodes2Text()
            plain_text = converter.latex_to_text(latex_content)
            
            return self._clean_text_for_db(plain_text)
            
        except Exception as e:
            logger.error(f"Error extracting LaTeX text: {str(e)}")
            # Fallback to plain text
            try:
                return await self._extract_text_file(file_path)
            except:
                return "Error extracting LaTeX content"
    
    async def _extract_rtf_text(self, file_path: str) -> str:
        """Extract text from RTF files"""
        try:
            # Try with striprtf library first (more reliable)
            try:
                from striprtf.striprtf import rtf_to_text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return self._clean_text_for_db(text)
            except ImportError:
                pass
                
            # Skip textract dependency - using basic parsing only
            
            # Basic RTF parsing as last resort
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            # Simple RTF text extraction (basic approach)
            import re
            # Remove RTF control codes
            text = re.sub(r'\\[a-z]+\d*\s?', ' ', content)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\s+', ' ', text).strip()
            
            return self._clean_text_for_db(text)
                
        except Exception as e:
            logger.error(f"Error extracting RTF text: {str(e)}")
            return "Error extracting RTF content"
    
    async def _extract_odt_text(self, file_path: str) -> str:
        """Extract text from OpenDocument files"""
        try:
            # Try with odfpy library first (more reliable)
            try:
                from odf.opendocument import load
                from odf.text import P
                from odf import teletype
                
                doc = load(file_path)
                text_content = []
                
                for paragraph in doc.getElementsByType(P):
                    text_content.append(teletype.extractText(paragraph))
                
                extracted_text = '\n'.join(text_content)
                return self._clean_text_for_db(extracted_text)
                
            except ImportError:
                pass
            
            # Skip textract dependency - using basic parsing only
            
            # Basic ODT parsing - ODT files are ZIP archives containing XML
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Try to read content.xml
                try:
                    content_xml = zip_file.read('content.xml')
                    root = ET.fromstring(content_xml)
                    
                    # Extract text from all text nodes
                    text_content = []
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            text_content.append(elem.text.strip())
                        if elem.tail and elem.tail.strip():
                            text_content.append(elem.tail.strip())
                    
                    extracted_text = ' '.join(text_content)
                    return self._clean_text_for_db(extracted_text)
                    
                except Exception as xml_error:
                    logger.warning(f"XML parsing failed: {xml_error}")
                    return "OpenDocument content could not be extracted - XML parsing failed"
                
        except Exception as e:
            logger.error(f"Error extracting OpenDocument text: {str(e)}")
            return "Error extracting OpenDocument content"
    
    async def _analyze_image_with_openai(self, image_base64: str, api_key: str) -> Dict[str, Any]:
        """Analyze image using OpenAI Vision API"""
        try:
            import openai
            client = openai.AsyncOpenAI(api_key=api_key)
            
            response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text", 
                                "text": "Analyze this image in detail. Describe what you see including objects, text, colors, and any important information. If there is text in the image, extract it."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{image_base64}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=1000
            )
            
            content = response.choices[0].message.content
            
            # Generate summary
            summary_response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": f"Summarize the main content of the image based on this analysis: {content}"
                    }
                ],
                max_tokens=200
            )
            
            return {
                "description": content,
                "analysis": content,
                "summary": summary_response.choices[0].message.content,
                "metadata": {
                    "model": "gpt-4o",
                    "provider": "openai"
                }
            }
            
        except Exception as e:
            logger.error(f"Error analyzing image with OpenAI: {str(e)}")
            return {
                "description": "Error analyzing image",
                "analysis": f"Error: {str(e)}",
                "summary": "Could not analyze image",
                "metadata": {}
            }
    
    async def _analyze_image_with_claude(self, image_base64: str, api_key: str) -> Dict[str, Any]:
        """Analyze image using Claude Vision API"""
        try:
            import anthropic
            client = anthropic.AsyncAnthropic(api_key=api_key)
            
            response = await client.messages.create(
                model="claude-3-sonnet-20240229",
                max_tokens=1000,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": "image/jpeg",
                                    "data": image_base64
                                }
                            },
                            {
                                "type": "text",
                                "text": "Analyze this image in detail. Describe what you see including objects, text, colors, and any important information. If there is text in the image, extract it."
                            }
                        ]
                    }
                ]
            )
            
            content = response.content[0].text
            
            # Generate summary  
            summary_response = await client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=200,
                messages=[
                    {
                        "role": "user",
                        "content": f"Summarize the main content of the image based on this analysis: {content}"
                    }
                ]
            )
            
            return {
                "description": content,
                "analysis": content,
                "summary": summary_response.content[0].text,
                "metadata": {
                    "model": "claude-3-sonnet-20240229",
                    "provider": "anthropic"
                }
            }
            
        except Exception as e:
            logger.error(f"Error analyzing image with Claude: {str(e)}")
            return {
                "description": "Error analyzing image",
                "analysis": f"Error: {str(e)}",
                "summary": "Could not analyze image",
                "metadata": {}
            }
    
    async def _analyze_image_with_gemini(self, image_base64: str, api_key: str) -> Dict[str, Any]:
        """Analyze image using official Google AI Client API - following image understanding docs"""
        try:
            from google import genai
            from google.genai import types
            
            # Create client with API key
            client = genai.Client(api_key=api_key)
            
            # Convert base64 to bytes
            image_data = base64.b64decode(image_base64)
            
            # Detect image format for correct MIME type
            mime_type = 'image/jpeg'  # Default
            if image_data.startswith(b'\x89PNG'):
                mime_type = 'image/png'
            elif image_data.startswith(b'GIF'):
                mime_type = 'image/gif'
            elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:20]:
                mime_type = 'image/webp'
            
            # Analysis prompt following Google's best practices
            analysis_prompt = """Analyze this image comprehensively and provide:

1. **Visual Content**: Describe what you see in detail
2. **Text Extraction (OCR)**: Extract any text visible in the image
3. **Objects & People**: Identify and describe objects, people, or subjects
4. **Visual Elements**: Colors, composition, style, quality
5. **Context & Purpose**: Infer the context and purpose of this image

Provide detailed analysis in English."""
            
            # Generate analysis using official API format
            def generate_analysis():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[
                        types.Part.from_bytes(
                            data=image_data,
                            mime_type=mime_type
                        ),
                        analysis_prompt
                    ]
                )
            
            response = await asyncio.get_event_loop().run_in_executor(None, generate_analysis)
            content = response.text
            
            # Generate concise summary
            summary_prompt = f"Provide a concise summary of the main content in this image (max 3 sentences): {content[:1000]}"
            
            def generate_summary():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[summary_prompt]
                )
            
            summary_response = await asyncio.get_event_loop().run_in_executor(None, generate_summary)
            
            return {
                "description": content,
                "analysis": content,
                "summary": summary_response.text,
                "metadata": {
                    "model": "gemini-2.5-flash",
                    "provider": "google",
                    "mime_type": mime_type,
                    "api_version": "google_ai_client",
                    "processing_method": "vision_understanding"
                }
            }
            
        except Exception as e:
            logger.error(f"Error analyzing image with Gemini: {str(e)}")
            return {
                "description": "Error analyzing image",
                "analysis": f"Gemini image analysis error: {str(e)}",
                "summary": "Could not analyze image with Gemini",
                "metadata": {"error": str(e)}
            }
    
    async def _process_video(self, file_path: str, ai_provider: str, api_key: str) -> Dict[str, Any]:
        """Process video using Gemini API according to official docs"""
        try:
            if ai_provider == "gemini" and api_key:
                return await self._analyze_video_with_gemini(file_path, api_key)
            else:
                # Basic video info without AI
                file_size = os.path.getsize(file_path)
                return {
                    "extracted_text": f"Video file uploaded: {os.path.basename(file_path)} ({file_size} bytes)",
                    "analysis": "Video processing requires Gemini API key for AI-powered analysis. Upload successful.",
                    "summary": "Video uploaded - AI analysis available with Gemini API",
                    "metadata": {"type": "video", "file_size": file_size, "status": "uploaded_only"}
                }
                
        except Exception as e:
            logger.error(f"Error processing video {file_path}: {str(e)}")
            raise
    
    async def _process_audio(self, file_path: str, ai_provider: str, api_key: str) -> Dict[str, Any]:
        """Process audio using Gemini API according to official docs"""
        try:
            if ai_provider == "gemini" and api_key:
                return await self._analyze_audio_with_gemini(file_path, api_key)
            else:
                # Basic audio info without AI
                file_size = os.path.getsize(file_path)
                return {
                    "extracted_text": f"Audio file uploaded: {os.path.basename(file_path)} ({file_size} bytes)",
                    "analysis": "Audio processing requires Gemini API key for speech recognition and analysis. Upload successful.",
                    "summary": "Audio uploaded - AI transcription available with Gemini API",
                    "metadata": {"type": "audio", "file_size": file_size, "status": "uploaded_only"}
                }
                
        except Exception as e:
            logger.error(f"Error processing audio {file_path}: {str(e)}")
            raise
    
    async def _analyze_video_with_gemini(self, file_path: str, api_key: str) -> Dict[str, Any]:
        """Analyze video using official Google AI video understanding capabilities"""
        try:
            from google import genai
            from google.genai import types
            
            # Create client with API key
            client = genai.Client(api_key=api_key)
            
            # Read video file as bytes
            with open(file_path, 'rb') as f:
                video_data = f.read()
            
            # Get file info
            file_size_mb = len(video_data) / (1024 * 1024)
            
            # Check file size limits (reasonable for processing)
            if file_size_mb > 100:  # Conservative limit
                return {
                    "extracted_text": f"Video file too large: {file_size_mb:.1f}MB",
                    "analysis": f"Video file ({file_size_mb:.1f}MB) exceeds processing limit. Consider a smaller file or use video trimming.",
                    "summary": "Video too large for direct processing",
                    "metadata": {"file_size_mb": file_size_mb, "status": "too_large"}
                }
            
            # Detect video MIME type
            mime_type = 'video/mp4'  # Default
            if file_path.lower().endswith('.mov'):
                mime_type = 'video/mov'
            elif file_path.lower().endswith('.avi'):
                mime_type = 'video/avi'
            elif file_path.lower().endswith('.webm'):
                mime_type = 'video/webm'
            elif file_path.lower().endswith('.3gpp'):
                mime_type = 'video/3gpp'
            
            # Video analysis prompt following official docs
            analysis_prompt = """Analyze this video comprehensively using your video understanding capabilities:

1. **Visual Content**: Describe what you see throughout the video
2. **Audio Content**: Describe any audio, speech, music, or sounds
3. **Scene Analysis**: Break down key scenes and their content
4. **Objects & People**: Identify objects, people, and activities
5. **Temporal Understanding**: Describe how content changes over time
6. **Text Recognition**: Extract any visible text or captions
7. **Context & Purpose**: Understand the video's purpose and context

Please provide detailed analysis in English, utilizing your video understanding and temporal reasoning capabilities."""
            
            # Generate analysis using official video API
            def generate_analysis():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[
                        types.Part.from_bytes(
                            data=video_data,
                            mime_type=mime_type
                        ),
                        analysis_prompt
                    ]
                )
            
            response = await asyncio.get_event_loop().run_in_executor(None, generate_analysis)
            content = response.text
            
            # Generate summary
            summary_prompt = f"Summarize the main content of this video (max 3 sentences): {content[:1000]}"
            
            def generate_summary():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[summary_prompt]
                )
            
            summary_response = await asyncio.get_event_loop().run_in_executor(None, generate_summary)
            
            return {
                "extracted_text": self._clean_text_for_db(content),
                "analysis": self._clean_text_for_db(content),
                "summary": self._clean_text_for_db(summary_response.text),
                "metadata": {
                    "model": "gemini-2.5-flash",
                    "provider": "google",
                    "mime_type": mime_type,
                    "file_size_mb": round(file_size_mb, 2),
                    "processing_method": "video_understanding",
                    "api_version": "google_ai_client",
                    "capabilities": "scene_analysis,temporal_reasoning,audio_visual"
                }
            }
            
        except Exception as e:
            logger.error(f"Error in video analysis: {str(e)}")
            # Fallback to basic info
            file_size = os.path.getsize(file_path)
            return {
                "extracted_text": f"Video file: {os.path.basename(file_path)} ({file_size} bytes)",
                "analysis": f"Video analysis error: {str(e)}. File uploaded successfully for manual review.",
                "summary": "Video processing encountered an error but file is stored",
                "metadata": {"error": str(e), "file_size": file_size, "fallback": True}
            }
    
    async def _analyze_audio_with_gemini(self, file_path: str, api_key: str) -> Dict[str, Any]:
        """Analyze audio using official Google AI audio understanding capabilities"""
        try:
            from google import genai
            from google.genai import types
            
            # Create client with API key
            client = genai.Client(api_key=api_key)
            
            # Read audio file as bytes
            with open(file_path, 'rb') as f:
                audio_data = f.read()
            
            # Get file info
            file_size_mb = len(audio_data) / (1024 * 1024)
            
            # Check file size limits
            if file_size_mb > 50:  # Conservative limit for audio
                return {
                    "extracted_text": f"Audio file too large: {file_size_mb:.1f}MB",
                    "analysis": f"Audio file ({file_size_mb:.1f}MB) exceeds processing limit. Consider audio compression or trimming.",
                    "summary": "Audio too large for direct processing",
                    "metadata": {"file_size_mb": file_size_mb, "status": "too_large"}
                }
            
            # Detect audio MIME type
            mime_type = 'audio/wav'  # Default
            if file_path.lower().endswith('.mp3'):
                mime_type = 'audio/mp3'
            elif file_path.lower().endswith('.aac'):
                mime_type = 'audio/aac'
            elif file_path.lower().endswith('.ogg'):
                mime_type = 'audio/ogg'
            elif file_path.lower().endswith('.flac'):
                mime_type = 'audio/flac'
            elif file_path.lower().endswith('.aiff'):
                mime_type = 'audio/aiff'
            
            # Audio analysis prompt following official docs
            analysis_prompt = """Analyze this audio comprehensively using your audio understanding capabilities:

1. **Speech Recognition**: Transcribe any speech or spoken content
2. **Audio Content**: Describe music, sounds, or audio elements
3. **Speaker Analysis**: Identify different speakers or voices if present
4. **Audio Quality**: Assess audio quality, clarity, and technical aspects
5. **Context & Purpose**: Understand the audio's purpose and context
6. **Emotional Tone**: Analyze tone, mood, or emotional content
7. **Background Sounds**: Identify any background sounds or environment

Please provide detailed analysis in English, including full transcription of any speech content."""
            
            # Generate analysis using official audio API
            def generate_analysis():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[
                        types.Part.from_bytes(
                            data=audio_data,
                            mime_type=mime_type
                        ),
                        analysis_prompt
                    ]
                )
            
            response = await asyncio.get_event_loop().run_in_executor(None, generate_analysis)
            content = response.text
            
            # Generate summary
            summary_prompt = f"Summarize the main content of this audio (max 3 sentences): {content[:1000]}"
            
            def generate_summary():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[summary_prompt]
                )
            
            summary_response = await asyncio.get_event_loop().run_in_executor(None, generate_summary)
            
            return {
                "extracted_text": self._clean_text_for_db(content),
                "analysis": self._clean_text_for_db(content),
                "summary": self._clean_text_for_db(summary_response.text),
                "metadata": {
                    "model": "gemini-2.5-flash",
                    "provider": "google",
                    "mime_type": mime_type,
                    "file_size_mb": round(file_size_mb, 2),
                    "processing_method": "audio_understanding",
                    "api_version": "google_ai_client",
                    "capabilities": "speech_recognition,audio_analysis,transcription"
                }
            }
            
        except Exception as e:
            logger.error(f"Error in audio analysis: {str(e)}")
            # Fallback to basic info
            file_size = os.path.getsize(file_path)
            return {
                "extracted_text": f"Audio file: {os.path.basename(file_path)} ({file_size} bytes)",
                "analysis": f"Audio analysis error: {str(e)}. File uploaded successfully for manual review.",
                "summary": "Audio processing encountered an error but file is stored",
                "metadata": {"error": str(e), "file_size": file_size, "fallback": True}
            }
    
    async def _get_basic_video_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic video information without AI"""
        try:
            file_size = os.path.getsize(file_path)
            description = f"Video file ({file_size} bytes)"
            
            return {
                "extracted_text": description,
                "analysis": "Basic video information (AI analysis not available)",
                "summary": description,
                "metadata": {"file_size": file_size}
            }
        except Exception as e:
            return {
                "extracted_text": "Could not process video",
                "analysis": f"Error: {str(e)}",
                "summary": "Video processing failed",
                "metadata": {}
            }
    
    async def _get_basic_audio_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic audio information without AI"""
        try:
            file_size = os.path.getsize(file_path)
            description = f"Audio file ({file_size} bytes)"
            
            return {
                "extracted_text": description,
                "analysis": "Basic audio information (AI analysis not available)",
                "summary": description,
                "metadata": {"file_size": file_size}
            }
        except Exception as e:
            return {
                "extracted_text": "Could not process audio",
                "analysis": f"Error: {str(e)}",
                "summary": "Audio processing failed",
                "metadata": {}
            }

    async def _get_basic_image_info(self, file_path: str) -> Dict[str, Any]:
        """Get basic image information without AI"""
        try:
            if IMAGE_AVAILABLE:
                with Image.open(file_path) as img:
                    width, height = img.size
                    format_name = img.format
                    mode = img.mode
                    
                    description = f"Image: {format_name} format, {width}x{height} pixels, {mode} mode"
            else:
                file_size = os.path.getsize(file_path)
                description = f"Image file ({file_size} bytes)"
            
            return {
                "description": description,
                "analysis": "Basic image information (AI analysis not available)",
                "summary": description,
                "metadata": {}
            }
            
        except Exception as e:
            return {
                "description": "Could not process image",
                "analysis": f"Error: {str(e)}",
                "summary": "Image processing failed",
                "metadata": {}
            }
    
    async def _analyze_text_with_ai(self, text: str, ai_provider: str, api_key: str) -> Dict[str, Any]:
        """Analyze extracted text using AI"""
        try:
            if ai_provider == "openai":
                return await self._analyze_text_with_openai(text, api_key)
            elif ai_provider == "claude":
                return await self._analyze_text_with_claude(text, api_key)
            elif ai_provider == "gemini":
                return await self._analyze_text_with_gemini(text, api_key)
            else:
                return {
                    "analysis": "AI analysis not available for this provider",
                    "summary": text[:200] + "..." if len(text) > 200 else text
                }
                
        except Exception as e:
            logger.error(f"Error analyzing text with AI: {str(e)}")
            return {
                "analysis": f"Error analyzing text: {str(e)}",
                "summary": text[:200] + "..." if len(text) > 200 else text
            }
    
    async def _analyze_text_with_openai(self, text: str, api_key: str) -> Dict[str, Any]:
        """Analyze text using OpenAI"""
        try:
            import openai
            client = openai.AsyncOpenAI(api_key=api_key)
            
            # Analysis
            analysis_response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": f"Analyze the following document and provide detailed information about content, main topics, and important points:\n\n{text[:4000]}"
                    }
                ],
                max_tokens=1000
            )
            
            # Summary
            summary_response = await client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": f"Provide a concise summary of the main content of the following document:\n\n{text[:2000]}"
                    }
                ],
                max_tokens=300
            )
            
            return {
                "analysis": analysis_response.choices[0].message.content,
                "summary": summary_response.choices[0].message.content
            }
            
        except Exception as e:
            raise Exception(f"OpenAI analysis error: {str(e)}")
    
    async def _analyze_text_with_claude(self, text: str, api_key: str) -> Dict[str, Any]:
        """Analyze text using Claude"""
        try:
            import anthropic
            client = anthropic.AsyncAnthropic(api_key=api_key)
            
            # Analysis
            analysis_response = await client.messages.create(
                model="claude-3-sonnet-20240229",
                max_tokens=1000,
                messages=[
                    {
                        "role": "user",
                        "content": f"Analyze the following document and provide detailed information about content, main topics, and important points:\n\n{text[:4000]}"
                    }
                ]
            )
            
            # Summary
            summary_response = await client.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=300,
                messages=[
                    {
                        "role": "user",
                        "content": f"Provide a concise summary of the main content of the following document:\n\n{text[:2000]}"
                    }
                ]
            )
            
            return {
                "analysis": analysis_response.content[0].text,
                "summary": summary_response.content[0].text
            }
            
        except Exception as e:
            raise Exception(f"Claude analysis error: {str(e)}")
    
    async def _analyze_text_with_gemini(self, text: str, api_key: str) -> Dict[str, Any]:
        """Analyze text using new Google Gemini Client API"""
        try:
            from google import genai
            from google.genai import types
            
            # Create client with API key
            client = genai.Client(api_key=api_key)
            
            # Analysis prompt in English
            analysis_prompt = f"""Analyze the following document and provide detailed information about:
- Main content and topics
- Key points and important information
- Document structure and organization
- Purpose and target audience

Document:
{text[:4000]}

Please respond in English."""
            
            # Generate analysis using new API
            def generate_analysis():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[analysis_prompt]
                )
            
            analysis_response = await asyncio.get_event_loop().run_in_executor(None, generate_analysis)
            
            # Summary
            summary_prompt = f"Summarize the main content of the document in English (max 200 words):\n\n{text[:2000]}"
            
            def generate_summary():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[summary_prompt]
                )
            
            summary_response = await asyncio.get_event_loop().run_in_executor(None, generate_summary)
            
            return {
                "analysis": analysis_response.text,
                "summary": summary_response.text
            }
            
        except Exception as e:
            raise Exception(f"Gemini analysis error: {str(e)}")
    
    async def _process_pdf_with_gemini(self, file_path: str, ai_provider: str, api_key: str) -> Dict[str, Any]:
        """Process PDF using official Google AI document processing API"""
        try:
            from google import genai
            from google.genai import types
            
            # Create client with API key
            client = genai.Client(api_key=api_key)
            
            # Read PDF file as bytes
            with open(file_path, 'rb') as f:
                pdf_data = f.read()
            
            # Check file size - Gemini supports up to 1000 pages (each page = 258 tokens)
            file_size_mb = len(pdf_data) / (1024 * 1024)
            if file_size_mb > 200:  # Reasonable limit for processing
                raise Exception(f"PDF file too large: {file_size_mb:.1f}MB. Maximum supported: ~200MB")
            
            # Document analysis prompt following official docs
            analysis_prompt = """Analyze this PDF document comprehensively using your document understanding capabilities:

1. **Document Overview**: What type of document is this and what's its main purpose?
2. **Content Analysis**: Analyze text, images, diagrams, charts, and tables throughout the document
3. **Key Information**: Extract the most important information and insights
4. **Structure & Organization**: How is the document structured and organized?
5. **Text Extraction**: Extract all readable text content accurately
6. **Visual Elements**: Describe any images, charts, diagrams, or visual content

Please provide a thorough analysis in Vietnamese, utilizing your natural vision capabilities to understand the complete context of this document."""
            
            # Generate analysis using official document processing API
            def generate_analysis():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[
                        types.Part.from_bytes(
                            data=pdf_data,
                            mime_type='application/pdf'
                        ),
                        analysis_prompt
                    ]
                )
            
            response = await asyncio.get_event_loop().run_in_executor(None, generate_analysis)
            content = response.text
            
            # Generate executive summary
            summary_prompt = f"Create a concise executive summary of this PDF document (max 200 words): {content[:2000]}"
            
            def generate_summary():
                return client.models.generate_content(
                    model="gemini-2.5-flash",
                    contents=[summary_prompt]
                )
            
            summary_response = await asyncio.get_event_loop().run_in_executor(None, generate_summary)
            
            return {
                "extracted_text": self._clean_text_for_db(content),
                "analysis": self._clean_text_for_db(content),
                "summary": self._clean_text_for_db(summary_response.text),
                "metadata": {
                    "model": "gemini-2.5-flash",
                    "provider": "google",
                    "processing_method": "document_understanding",
                    "api_version": "google_ai_client",
                    "file_size_mb": round(file_size_mb, 2),
                    "supports_up_to": "1000_pages_max",
                    "token_estimate": f"~{len(pdf_data) // 1000}_tokens"
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF with Gemini: {str(e)}")
            raise Exception(f"Gemini PDF processing error: {str(e)}")
